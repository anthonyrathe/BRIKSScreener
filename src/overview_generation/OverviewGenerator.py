import os, datetime
from os.path import dirname as dirname
from pptx import Presentation
from src.data_processing.DataReader import DataReader
from src.exceptions.NoDataFoundException import NoDataFoundException
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.dml.color import RGBColor
from pptx.chart.data import ChartData

def clip(str,n=1500):
			if len(str)>n:
				return str[:n-3]+"..."
			return str

class OverviewGenerator:
	def __init__(self):
		self.master_path = os.path.relpath("{}/master.pptx".format(dirname(__file__)))
		self.master = Presentation(self.master_path)
		self.slide_index = 0

	def generateOverview(self,ticker,update=[]):
		data = DataReader(ticker).getOverview(update=update)

		slide = self.master.slides[self.slide_index]

		slide.shapes[0].text = "Company Overview Generator ({})".format(datetime.datetime.now().date().strftime(format="%d-%m-%Y"))
		slide.shapes[3].text = "{} ({})".format(data['name'],data['ticker'])
		slide.shapes[9].text = clip(data['summary'])
		slide.shapes[9].text_frame.paragraphs[0].runs[0].font.size = Pt(8)
		slide.shapes[9].text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
		slide.shapes[8].text = "*{}".format(data['industry'])
		slide.shapes[8].text_frame.paragraphs[0].runs[0].font.size = Pt(9)
		slide.shapes[8].text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT

		def to_multiple(f):
			return "{}x".format(round(f,2))

		def to_percentage(f):
			return "{}%".format(round(f*100,2))

		def to_formatted_amount(f,curr):
			if f >= 10**12:
				return "{}T {}".format(round(f/(10**12),2),curr)
			elif f >= 10**9:
				return "{}B {}".format(round(f/(10**9),2),curr)
			elif f >= 10**6:
				return "{}M {}".format(round(f/(10**6),2),curr)
			elif f >= 10**3:
				return "{}K {}".format(round(f/(10**3),2),curr)

		table_contents = [[data['ticker'],"Industry Median*"],]
		table_contents.append([to_multiple(data['multiples']['EV/EBITDA']['company']),to_multiple(data['multiples']['EV/EBITDA']['peers'])])
		table_contents.append([to_multiple(data['multiples']['P/E']['company']),to_multiple(data['multiples']['P/E']['peers'])])
		table_contents.append([to_multiple(data['multiples']['P/B']['company']),to_multiple(data['multiples']['P/B']['peers'])])
		table_contents.append([to_percentage(data['multiples']['ROE']['company']),to_percentage(data['multiples']['ROE']['peers'])])
		table_contents.append([to_multiple(data['multiples']['current_ratio']['company']),to_multiple(data['multiples']['current_ratio']['peers'])])
		table_contents.append([to_percentage(data['multiples']['dividend_yield']['company']),to_percentage(data['multiples']['dividend_yield']['peers'])])
		table_contents.append([to_formatted_amount(data['multiples']['market_cap']['company'],data['currencies']['exchange']),to_formatted_amount(data['multiples']['market_cap']['peers'],data['currencies']['exchange'])])

		for r in range(len(table_contents)):
			for i in range(len(table_contents[r])):
				slide.shapes[7].table.cell(r,1+i).text = str(table_contents[r][i])
				for paragraph in slide.shapes[7].table.cell(r,1+i).text_frame.paragraphs:
					paragraph.alignment = PP_ALIGN.CENTER
					for run in paragraph.runs:
						run.font.size = Pt(12)


		performance = data['performance'].reindex(index=data['performance'].index[::-1])

		# Add revenue chart
		shape_num = 10
		slide.shapes[shape_num].chart.chart_title.text_frame.text = "Revenue ({})".format(data['currencies']['filings'])

		chart_data = ChartData()
		chart_data.categories = list(performance['labels'])
		if abs(performance['totalrevenue']).max()>10**12:
			performance['totalrevenue'] /= 10**12
			slide.shapes[shape_num].chart.value_axis.axis_title.text_frame.text = "Trillions"
			slide.shapes[shape_num].chart.value_axis.axis_title.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor.from_string('595959')
			slide.shapes[shape_num].chart.value_axis.axis_title.text_frame.paragraphs[0].runs[0].font.size = Pt(9)
		elif abs(performance['totalrevenue']).max()>10**9:
			performance['totalrevenue'] /= 10**9
			slide.shapes[shape_num].chart.value_axis.axis_title.text_frame.text = "Billions"
			slide.shapes[shape_num].chart.value_axis.axis_title.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor.from_string('595959')
			slide.shapes[shape_num].chart.value_axis.axis_title.text_frame.paragraphs[0].runs[0].font.size = Pt(9)
		elif abs(performance['totalrevenue']).max()>10**6:
			performance['totalrevenue'] /= 10**6
			slide.shapes[shape_num].chart.value_axis.axis_title.text_frame.text = "Millions"
			slide.shapes[shape_num].chart.value_axis.axis_title.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor.from_string('595959')
			slide.shapes[shape_num].chart.value_axis.axis_title.text_frame.paragraphs[0].runs[0].font.size = Pt(9)
		elif abs(performance['totalrevenue']).max()>10**3:
			performance['totalrevenue'] /= 10**3
			slide.shapes[shape_num].chart.value_axis.axis_title.text_frame.text = "Thousands"
			slide.shapes[shape_num].chart.value_axis.axis_title.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor.from_string('595959')
			slide.shapes[shape_num].chart.value_axis.axis_title.text_frame.paragraphs[0].runs[0].font.size = Pt(9)
		chart_data.add_series('Revenue', tuple(performance['totalrevenue']))
		slide.shapes[shape_num].chart.replace_data(chart_data)



		# Add EBITDA & earnings chart
		shape_num = 11
		slide.shapes[shape_num].chart.chart_title.text_frame.text = "EBITDA & Net Earnings ({})".format(data['currencies']['filings'])

		chart_data = ChartData()
		chart_data.categories = list(performance['labels'])
		if max(abs(performance['EBITDA']).max(),abs(performance['netincome']).max())>10**12:
			performance['EBITDA'] /= 10**12
			performance['netincome'] /= 10**12
			slide.shapes[shape_num].chart.value_axis.axis_title.text_frame.text = "Trillions"
			slide.shapes[shape_num].chart.value_axis.axis_title.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor.from_string('595959')
			slide.shapes[shape_num].chart.value_axis.axis_title.text_frame.paragraphs[0].runs[0].font.size = Pt(9)
		elif max(abs(performance['EBITDA']).max(),abs(performance['netincome']).max())>10**9:
			performance['EBITDA'] /= 10**9
			performance['netincome'] /= 10**9
			slide.shapes[shape_num].chart.value_axis.axis_title.text_frame.text = "Billions"
			slide.shapes[shape_num].chart.value_axis.axis_title.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor.from_string('595959')
			slide.shapes[shape_num].chart.value_axis.axis_title.text_frame.paragraphs[0].runs[0].font.size = Pt(9)
		elif max(abs(performance['EBITDA']).max(),abs(performance['netincome']).max())>10**6:
			performance['EBITDA'] /= 10**6
			performance['netincome'] /= 10**6
			slide.shapes[shape_num].chart.value_axis.axis_title.text_frame.text = "Millions"
			slide.shapes[shape_num].chart.value_axis.axis_title.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor.from_string('595959')
			slide.shapes[shape_num].chart.value_axis.axis_title.text_frame.paragraphs[0].runs[0].font.size = Pt(9)
		elif max(abs(performance['EBITDA']).max(),abs(performance['netincome']).max())>10**3:
			performance['EBITDA'] /= 10**3
			performance['netincome'] /= 10**3
			slide.shapes[shape_num].chart.value_axis.axis_title.text_frame.text = "Thousands"
			slide.shapes[shape_num].chart.value_axis.axis_title.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor.from_string('595959')
			slide.shapes[shape_num].chart.value_axis.axis_title.text_frame.paragraphs[0].runs[0].font.size = Pt(9)
		chart_data.add_series('EBITDA', tuple(performance['EBITDA']))
		chart_data.add_series('Net Earnings', tuple(performance['netincome']))
		slide.shapes[shape_num].chart.replace_data(chart_data)

		# Add prices chart
		shape_num = 12
		currency = data['currencies']['exchange']
		if currency == 'GBP':
			currency = 'GBp'
		slide.shapes[shape_num].chart.chart_title.text_frame.text = "Price ({})".format(currency)

		prices = data['prices']
		chart_data = ChartData()
		chart_data.categories = list(prices.index)
		chart_data.add_series('Revenue', tuple(prices))
		try:
			slide.shapes[shape_num].chart.replace_data(chart_data)
		except TypeError:
			raise NoDataFoundException


		for i in [10,11,12]:
			slide.shapes[i].chart.chart_title.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor.from_string('595959')
			slide.shapes[i].chart.chart_title.text_frame.paragraphs[0].runs[0].font.size = Pt(14)

		self.slide_index += 1

		return self.slide_index < len(self.master.slides)

	def save(self,filename):
		if self.slide_index < len(self.master.slides):
			for i in range(len(self.master.slides)-1,self.slide_index-1,-1):
				rId = self.master.slides._sldIdLst[i].rId
				self.master.part.drop_rel(rId)
				del self.master.slides._sldIdLst[i]

		self.master.save(os.path.relpath("{}/data/cleaned/overview_ppt/{}.pptx".format(dirname(dirname(dirname(__file__))),filename)))
